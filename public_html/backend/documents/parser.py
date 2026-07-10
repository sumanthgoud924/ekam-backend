import io
import re
import logging
from pathlib import Path
from typing import Optional, BinaryIO

logger = logging.getLogger(__name__)


class DocumentParser:
    SUPPORTED_EXTENSIONS = {
        ".pdf", ".docx", ".doc", ".txt", ".pptx", ".xlsx",
        ".epub", ".md", ".rtf", ".html", ".htm", ".odt",
        ".csv", ".json", ".xml", ".png", ".jpg", ".jpeg",
    }

    @staticmethod
    def detect_language(text: str) -> str:
        try:
            import langdetect
            return langdetect.detect(text)
        except Exception:
            return "en"

    @staticmethod
    def clean_text(text: str) -> str:
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r'[ \t]+$', '', text, flags=re.MULTILINE)
        return text.strip()

    async def parse(self, file_path: str | Path) -> dict:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")
        ext = path.suffix.lower()

        parsers = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".doc": self._parse_docx,
            ".txt": self._parse_text,
            ".pptx": self._parse_pptx,
            ".xlsx": self._parse_xlsx,
            ".epub": self._parse_epub,
            ".md": self._parse_text,
            ".rtf": self._parse_text,
            ".html": self._parse_html,
            ".htm": self._parse_html,
            ".odt": self._parse_docx,
            ".csv": self._parse_text,
            ".json": self._parse_text,
            ".xml": self._parse_text,
        }

        parser = parsers.get(ext)
        if parser is None:
            if ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp"}:
                return await self._parse_image(file_path)
            raise ValueError(f"Unsupported format: {ext}")

        return await parser(path)

    async def _parse_pdf(self, path: Path) -> dict:
        try:
            import fitz
        except ImportError:
            return await self._parse_pdf_fallback(path)

        doc = fitz.open(path)
        pages = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            pages.append({
                "page": page_num + 1,
                "text": text,
                "has_images": len(page.get_images()) > 0,
            })
        doc.close()

        full_text = "\n\n".join(p["text"] for p in pages)
        full_text = self.clean_text(full_text)
        language = self.detect_language(full_text)

        return {
            "content": full_text,
            "pages": len(pages),
            "page_details": pages,
            "total_characters": len(full_text),
            "language": language,
            "formats_detected": ["pdf"],
        }

    async def _parse_pdf_fallback(self, path: Path) -> dict:
        try:
            import subprocess
            result = subprocess.run(
                ["pdftotext", str(path), "-"],
                capture_output=True, text=True, check=False
            )
            full_text = result.stdout
        except Exception:
            with open(path, "rb") as f:
                raw = f.read()
            text = raw.decode("latin-1")
            text = re.sub(r'<</[^>]*>>', '', text)
            text = re.sub(r'stream.*?endstream', '', text, flags=re.DOTALL)
            full_text = re.sub(r'[^a-zA-Z0-9\s.,!?;:\'\"-]', ' ', text)

        full_text = self.clean_text(full_text)
        return {
            "content": full_text,
            "pages": full_text.count("\f") + 1,
            "page_details": [],
            "total_characters": len(full_text),
            "language": self.detect_language(full_text),
            "formats_detected": ["pdf"],
        }

    async def _parse_docx(self, path: Path) -> dict:
        try:
            from docx import Document
            doc = Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)
            full_text = self.clean_text(full_text)

            return {
                "content": full_text,
                "pages": 1,
                "page_details": [],
                "total_characters": len(full_text),
                "language": self.detect_language(full_text),
                "formats_detected": ["docx"],
            }
        except Exception as e:
            raise ValueError(f"Failed to parse DOCX: {e}")

    async def _parse_pptx(self, path: Path) -> dict:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text for cell in row.cells]
                            slide_text.append(" | ".join(row_text))
                slides.append({
                    "slide": i,
                    "text": "\n".join(slide_text),
                })

            full_text = "\n\n".join(s["text"] for s in slides)
            full_text = self.clean_text(full_text)

            return {
                "content": full_text,
                "pages": len(slides),
                "page_details": slides,
                "total_characters": len(full_text),
                "language": self.detect_language(full_text),
                "formats_detected": ["pptx"],
            }
        except Exception as e:
            raise ValueError(f"Failed to parse PPTX: {e}")

    async def _parse_xlsx(self, path: Path) -> dict:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            sheets_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in row_text):
                        rows.append("\t".join(row_text))
                if rows:
                    sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            wb.close()

            full_text = "\n\n".join(sheets_text)
            full_text = self.clean_text(full_text)

            return {
                "content": full_text,
                "pages": len(wb.sheetnames),
                "page_details": [],
                "total_characters": len(full_text),
                "language": self.detect_language(full_text),
                "formats_detected": ["xlsx"],
            }
        except Exception as e:
            raise ValueError(f"Failed to parse XLSX: {e}")

    async def _parse_epub(self, path: Path) -> dict:
        try:
            import zipfile
            import xml.etree.ElementTree as ET

            chapters = []
            with zipfile.ZipFile(path) as z:
                for name in z.namelist():
                    if name.endswith(".xhtml") or name.endswith(".html") or name.endswith(".htm"):
                        content = z.read(name).decode("utf-8", errors="replace")
                        text = re.sub(r'<[^>]+>', ' ', content)
                        text = re.sub(r'\s+', ' ', text).strip()
                        if text:
                            chapters.append(text)

            full_text = "\n\n".join(chapters)
            full_text = self.clean_text(full_text)

            return {
                "content": full_text,
                "pages": len(chapters),
                "page_details": [],
                "total_characters": len(full_text),
                "language": self.detect_language(full_text),
                "formats_detected": ["epub"],
            }
        except Exception as e:
            return self._read_text_fallback(path)

    async def _parse_html(self, path: Path) -> dict:
        try:
            from html.parser import HTMLParser

            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.text_parts = []
                    self.skip_tags = {"script", "style", "noscript"}

                def handle_data(self, data):
                    self.text_parts.append(data)

            content = path.read_text("utf-8", errors="replace")
            parser = TextExtractor()
            parser.feed(content)
            full_text = " ".join(parser.text_parts)
            full_text = re.sub(r'\s+', ' ', full_text).strip()

            return {
                "content": full_text,
                "pages": 1,
                "page_details": [],
                "total_characters": len(full_text),
                "language": self.detect_language(full_text),
                "formats_detected": ["html"],
            }
        except Exception as e:
            return self._read_text_fallback(path)

    async def _parse_text(self, path: Path) -> dict:
        return self._read_text_fallback(path)

    def _read_text_fallback(self, path: Path) -> dict:
        try:
            content = path.read_text("utf-8", errors="replace")
        except UnicodeDecodeError:
            content = path.read_text("latin-1", errors="replace")

        full_text = self.clean_text(content)
        return {
            "content": full_text,
            "pages": full_text.count("\n\n\n") + 1,
            "page_details": [],
            "total_characters": len(full_text),
            "language": self.detect_language(full_text),
            "formats_detected": [path.suffix.lower().lstrip(".")],
        }

    async def parse_with_chapters(self, file_path: str | Path) -> list[dict]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        ext = path.suffix.lower()
        text_formats = {".md", ".txt", ".rtf", ".csv", ".json", ".xml", ".html", ".htm"}
        heading_pattern = re.compile(r'^(#{1,4})\s+(.+)$|^([A-Z][A-Z\s]{2,})$')

        if ext in text_formats:
            try:
                raw = path.read_text("utf-8", errors="replace")
            except UnicodeDecodeError:
                raw = path.read_text("latin-1", errors="replace")
            lines = raw.split("\n")
        else:
            result = await self.parse(path)
            raw = result.get("content", "")
            lines = raw.split(". ")

        chapters = []
        current_title = "Introduction"
        current_lines = []

        def flush():
            text = "\n".join(current_lines).strip() if ext in text_formats else ". ".join(current_lines).strip()
            if text:
                cleaned = self.clean_text(text)
                chapters.append({
                    "title": current_title,
                    "content": cleaned or text,
                    "page_number": None,
                    "word_count": len(cleaned.split()) if cleaned else 0,
                })

        for line in lines:
            m = heading_pattern.match(line.strip())
            if m:
                flush()
                current_title = (m.group(2) or m.group(3)).strip()
                current_lines = []
            else:
                current_lines.append(line)

        flush()
        return chapters

    async def _parse_image(self, path: Path) -> dict:
        try:
            from PIL import Image
            import pytesseract

            img = Image.open(path)
            text = pytesseract.image_to_string(img)
            full_text = self.clean_text(text)
            language = self.detect_language(full_text)

            return {
                "content": full_text,
                "pages": 1,
                "page_details": [{"page": 1, "text": full_text, "has_images": False}],
                "total_characters": len(full_text),
                "language": language,
                "formats_detected": ["ocr"],
            }
        except ImportError:
            raise ValueError("OCR requires pytesseract and Pillow. Install: pip install pytesseract Pillow")
        except Exception as e:
            raise ValueError(f"OCR failed: {e}")
